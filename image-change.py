#!/usr/bin/env python3

import hashlib
import sys
import openpyxl
import pptx




def hexstuff():
    prs = pptx.Presentation('hack.pptx')
    newImgFilename = "gray.jpg"
    img2 = pptx.parts.image.Image.from_file(newImgFilename)

    print(hashlib.sha224(prs.slides[0].shapes[2].image._blob).hexdigest())
    print(hashlib.sha224(img2._blob).hexdigest())

    prs.slides[0].shapes[2].image._blob = img2._blob
    print(hashlib.sha224(prs.slides[0].shapes[2].image._blob).hexdigest())

prs = pptx.Presentation('big05.pptx')

wb = openpyxl.Workbook()
ws = wb.active


slide_num = 0
for slide in prs.slides:
    slide_num = slide_num + 1
    # print(slide_num)
    for shape in slide.shapes:
        for relpart in shape.part.related_parts.values():
            if isinstance(relpart, pptx.parts.image.ImagePart):
                imgsize = sys.getsizeof(relpart.image.blob)
                print(f"Slide {slide_num:>3} shape  {shape.shape_id:>4}, partname {relpart.partname:>30}, size {imgsize:>11} hash {relpart.sha1:<4}")
                xl_row = [slide_num, relpart.partname, imgsize, relpart.sha1]
                ws.append(xl_row)

wb.save('images.xlsx')


        # if (shape.is_placeholder or True):
            # phf = shape.placeholder_format
            # print(f'Slide # {slide_num:>3} shape idx {shape.shape_id:>4} placeholder idx {phf.idx:>2} type {phf.type:>8}')