#!/usr/bin/env python3

import pptx
prs = pptx.Presentation('test.pptx')


for slide in prs.slides:
    slide_num = slide_num + 1
    # print(slide_num)
    for shape in slide.shapes:
        for relpart in shape.part.related_parts.values():
            if isinstance(relpart, pptx.parts.image.ImagePart):
                print(f"Slide {slide_num:>3} shape  {shape.shape_id:>4}, partname {relpart.partname:>30}, size {imgsize:>11} hash {relpart.sha1:<4}")
